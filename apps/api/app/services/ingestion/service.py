import logging
import os
import re

from app.providers.base import StorageProvider

logger = logging.getLogger(__name__)

_ACADEMIC_SECTION_PATTERNS = [
    (re.compile(r"^abstract$", re.I), "abstract"),
    (re.compile(r"^introduction$", re.I), "introduction"),
    (re.compile(r"^related\s+work", re.I), "related_work"),
    (re.compile(r"^background", re.I), "background"),
    (re.compile(r"^(method|approach|model|architecture|system|framework)", re.I), "method"),
    (re.compile(r"^(experiment|evaluation|result|analysis)", re.I), "results"),
    (re.compile(r"^(discussion|limitation)", re.I), "discussion"),
    (re.compile(r"^(conclusion|summary|future)", re.I), "conclusion"),
    (re.compile(r"^(reference|bibliography)", re.I), "references"),
    (re.compile(r"^(appendix|supplement)", re.I), "appendix"),
    (re.compile(r"^(acknowledgment|acknowledgement)", re.I), "acknowledgments"),
]


def _classify_section(text: str) -> str | None:
    """Match a title-like text against common academic paper section patterns."""
    cleaned = text.strip().rstrip(".").strip()
    cleaned = re.sub(r"^\d+\.?\s*", "", cleaned).strip()
    for pattern, label in _ACADEMIC_SECTION_PATTERNS:
        if pattern.match(cleaned):
            return label
    return None


def _is_paper_title(text: str, page_num: int, font_size: float) -> bool:
    """Heuristic: large font on page 1 is likely the paper title."""
    return page_num == 1 and font_size > 18 and len(text.split()) >= 3


class IngestionService:
    def __init__(self, storage: StorageProvider):
        self.storage = storage

    async def process_pptx(self, file_path: str, doc_id: str) -> list[dict]:
        """Extract content from PPTX file using python-pptx."""
        from pptx import Presentation

        prs = Presentation(file_path)
        fragments: list[dict] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            ref_prefix = f"slide_{slide_num:02d}"

            if slide.shapes.title and slide.shapes.title.text.strip():
                fragments.append({
                    "ref_key": ref_prefix,
                    "page_or_slide_number": slide_num,
                    "kind": "title",
                    "text": slide.shapes.title.text.strip(),
                })

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if hasattr(slide.shapes, "title") and shape == slide.shapes.title:
                    continue
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if not text:
                        continue
                    fragments.append({
                        "ref_key": f"{ref_prefix}_bullet_{para_idx}",
                        "page_or_slide_number": slide_num,
                        "kind": "bullet",
                        "text": text,
                    })

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    fragments.append({
                        "ref_key": f"{ref_prefix}_note",
                        "page_or_slide_number": slide_num,
                        "kind": "note",
                        "text": notes_text,
                    })

        logger.info(
            "IngestionService: extracted %d fragments from PPTX doc_id=%s",
            len(fragments), doc_id,
        )
        return fragments

    async def process_pdf(self, file_path: str, doc_id: str) -> list[dict]:
        """Extract content from PDF, with academic paper section detection."""
        import fitz

        doc = fitz.open(file_path)
        raw_fragments: list[dict] = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            page_number = page_num + 1
            ref_prefix = f"page_{page_number:02d}"
            blocks = page.get_text("dict", sort=True).get("blocks", [])

            for block_idx, block in enumerate(blocks):
                if block.get("type") != 0:
                    continue

                block_texts: list[str] = []
                max_font_size = 0.0

                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        if text:
                            block_texts.append(text)
                            max_font_size = max(max_font_size, span.get("size", 12.0))

                combined_text = " ".join(block_texts)
                if not combined_text:
                    continue

                kind = "title" if max_font_size > 16 else "paragraph"

                academic_section = None
                if kind == "title":
                    academic_section = _classify_section(combined_text)
                    if _is_paper_title(combined_text, page_number, max_font_size):
                        academic_section = "paper_title"

                raw_fragments.append({
                    "ref_key": f"{ref_prefix}_block_{block_idx}",
                    "page_or_slide_number": page_number,
                    "kind": kind,
                    "text": combined_text,
                    "bbox": list(block.get("bbox", [])),
                    "font_size": max_font_size,
                    "academic_section": academic_section,
                })

        doc.close()

        # Merge consecutive paragraphs and attach section labels
        merged: list[dict] = []
        current: dict | None = None
        current_section: str | None = None

        for frag in raw_fragments:
            if frag.get("academic_section"):
                current_section = frag["academic_section"]

            frag["academic_section"] = frag.get("academic_section") or current_section

            if (
                current
                and current["kind"] == "paragraph"
                and frag["kind"] == "paragraph"
                and current["page_or_slide_number"] == frag["page_or_slide_number"]
            ):
                current["text"] += " " + frag["text"]
            else:
                if current:
                    merged.append(current)
                current = frag.copy()

        if current:
            merged.append(current)

        # Detect if this looks like an academic paper
        section_labels = {f.get("academic_section") for f in merged if f.get("academic_section")}
        is_academic = len(section_labels & {"abstract", "introduction", "conclusion", "method", "results"}) >= 2

        if is_academic:
            for frag in merged:
                if not frag.get("academic_section"):
                    frag["academic_section"] = "body"
            logger.info(
                "IngestionService: detected academic paper with sections: %s",
                sorted(section_labels),
            )

        logger.info(
            "IngestionService: extracted %d fragments from PDF doc_id=%s (academic=%s)",
            len(merged), doc_id, is_academic,
        )
        return merged

    async def process_topic(self, topic: str, domain: str) -> list[dict]:
        """Create synthetic source fragments from topic input."""
        logger.info("IngestionService: creating synthetic fragment for topic=%r", topic)
        return [{
            "ref_key": "synthetic_topic",
            "page_or_slide_number": None,
            "kind": "synthetic",
            "text": (
                f"Create a comprehensive educational lesson about: {topic}. "
                f"Domain: {domain}. "
                f"Cover all key concepts, definitions, algorithms, data structures, "
                f"mechanisms, and practical applications related to {topic}. "
                f"Include prerequisite knowledge, common misconceptions, "
                f"and worked examples where applicable."
            ),
        }]

    async def extract_fragments(
        self,
        source_type: str,
        file_path: str | None,
        topic: str | None,
        domain: str,
        doc_id: str,
    ) -> list[dict]:
        """Route to appropriate processor based on source type."""
        if source_type == "pptx" and file_path:
            if not os.path.isfile(file_path):
                logger.warning(
                    "IngestionService: PPTX not found, returning no fragments path=%r doc_id=%s",
                    file_path,
                    doc_id,
                )
                return []
            return await self.process_pptx(file_path, doc_id)
        elif source_type == "pdf" and file_path:
            if not os.path.isfile(file_path):
                logger.warning(
                    "IngestionService: PDF not found, returning no fragments path=%r doc_id=%s",
                    file_path,
                    doc_id,
                )
                return []
            return await self.process_pdf(file_path, doc_id)
        elif source_type == "topic" and topic:
            return await self.process_topic(topic, domain)
        else:
            raise ValueError(
                f"Unsupported source_type='{source_type}' or missing required input"
            )
